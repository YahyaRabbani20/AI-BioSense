# -*- coding: utf-8 -*-
"""
Created on Mon May  8 09:33:58 2023

@author: yahya
"""
                # # # # # 660 # # # # #
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

                            ### Load Data ###

# Load the data from a CSV file
df = pd.read_csv('ReversibilityBeforeFiltrationafter660 dec.csv', delimiter=',', header=None)


## Calculate the means and stds
time_point=1

totl_sample = 1
Analyte = 2
replicate = 3
Total = totl_sample * Analyte

total=df.shape[1]//time_point  # ensure total is an integer

# create lists for the means and stds
means = []
stds = []


# Loop for each time point
for t in range(time_point):
    start = t*total  # start of each time point

    time_point_means = []
    time_point_stds = []
    
    for i in range(Total):  # within each time point, step by 1
        indices = [start + i + j*Total for j in range(3)]
        # ensure all indices are within time point columns
        if all(idx < start + total for idx in indices):
            time_point_means.append(df.iloc[:, indices].mean(axis=1).tolist())
            time_point_stds.append(df.iloc[:, indices].std(axis=1).tolist())

    means.append(time_point_means)
    stds.append(time_point_stds)

# Now means and stds contain the means and standard deviations for each group of 96th columns (1, 97, 193, then 2, 98, 194, and so on) for each time point.

                           ## PLot the result  ###              

# Specify the time point and row to calculate the differences
## peak shift

# For Glucose, Mannose and Frouctose as analye
totl_sample = 1
Analyte = 2
replicate = 3
Total = totl_sample * Analyte

time_point = 0
peaks = [7,13]  # 
channels = ['(7,5)', '(7,6)']
x = np.arange(1, totl_sample+1)  # 1-indexing for DNA numbers
#labels for different analyte
labels = ['Glucose 7.5mM']
df = pd.DataFrame({'DNA Number': range(1, totl_sample + 1)})

fig, axs = plt.subplots(1, 2, figsize=(20, 5))  # only 2 subplots, one for each channel

for channel_idx, channel in enumerate(channels):
    axs[channel_idx].set_title(f'Channel {channel}')
    axs[channel_idx].set_xlabel('DNA Number')
    axs[channel_idx].set_ylabel('Peak Shift (nm)')
    # axs[channel_idx].set_ylim(-1, 1)  # Uncomment if you want to set the y-axis limits
    
    peak = peaks[channel_idx]
    for A in range(Analyte-1):
        means_diff = []
        stds_diff = []
        for DNA in range(totl_sample):
            idx = (A * totl_sample) + DNA
            blank_idx = ((Analyte - 1) * totl_sample) + DNA
            diff = means[time_point][idx][peak] - means[time_point][blank_idx][peak]
            err = np.sqrt(stds[time_point][idx][peak]**2 + stds[time_point][blank_idx][peak]**2)
            means_diff.append(diff)
            stds_diff.append(err)

        axs[channel_idx].bar(x + A*0.1, means_diff, yerr=stds_diff, width=0.1, align='center', alpha=0.5, capsize=5, label=labels[A])
        df[f'{channel} Shift'] = means_diff #####
        df[f'{channel} Shift_Std'] = stds_diff #####
        
    axs[channel_idx].set_xticks(x)
    axs[channel_idx].set_xticklabels(x)
    axs[channel_idx].legend()

# Save and display the plots
plt.tight_layout()
plt.savefig(f'660-All-Peak shift-t={time_point}.png', dpi=300, bbox_inches='tight')
plt.show()










# Intensity change

# peaks = [8,14] # 
# channels = ['(7,5)', '(7,6)']
# x = np.arange(1, totl_sample+1)  # 1-indexing for DNA numbers
# #labels for different analyte
# labels = ['Glucose 7.5mM']

# fig, axs = plt.subplots(1, 2, figsize=(20, 5))  # only 2 subplots, one for each channel

# for channel_idx, channel in enumerate(channels):
#     axs[channel_idx].set_title(f'Channel {channel}')
#     axs[channel_idx].set_xlabel('DNA Number')
#     axs[channel_idx].set_ylabel('Intensity change')
#     # axs[channel_idx].set_ylim(-1, 1)  # Uncomment if you want to set the y-axis limits
    
#     peak = peaks[channel_idx]
#     for A in range(Analyte-1):
#         means_diff = []
#         stds_diff = []
#         for DNA in range(totl_sample):
#             idx = (A * totl_sample) + DNA
#             blank_idx = ((Analyte - 1) * totl_sample) + DNA
#             diff = means[time_point][idx][peak] - means[time_point][blank_idx][peak]
#             err = np.sqrt(stds[time_point][idx][peak]**2 + stds[time_point][blank_idx][peak]**2)
#             means_diff.append(diff)
#             stds_diff.append(err)

#         axs[channel_idx].bar(x + A*0.1, means_diff, yerr=stds_diff, width=0.1, align='center', alpha=0.5, capsize=5, label=labels[A])

#     axs[channel_idx].set_xticks(x)
#     axs[channel_idx].set_xticklabels(x)
#     axs[channel_idx].legend()

# # Save and display the plots
# plt.tight_layout()
# plt.savefig(f'660-All-Intensity change-t={time_point}.png', dpi=300, bbox_inches='tight')
# plt.show()



#Normaliyed

peaks = [8,14] # 
channels = ['(7,5)', '(7,6)']
x = np.arange(1, totl_sample+1)  # 1-indexing for DNA numbers
#labels for different analyte
labels = ['Glucose 7.5mM']
df2 = pd.DataFrame({'DNA Number': range(1, totl_sample + 1)})

fig, axs = plt.subplots(1, 2, figsize=(20, 5))  # only 2 subplots, one for each channel

for channel_idx, channel in enumerate(channels):
    axs[channel_idx].set_title(f'Channel {channel}')
    axs[channel_idx].set_xlabel('DNA Number')
    axs[channel_idx].set_ylabel('Intensity change (%)')
    # axs[channel_idx].set_ylim(-1, 1)  # Uncomment if you want to set the y-axis limits
    
    peak = peaks[channel_idx]
    for A in range(Analyte-1):
        means_diff = []
        stds_diff = []
        for DNA in range(totl_sample):
            idx = (A * totl_sample) + DNA
            blank_idx = ((Analyte - 1) * totl_sample) + DNA
            diff = (means[time_point][idx][peak] - means[time_point][blank_idx][peak]) / means[time_point][blank_idx][peak] * 100
            err = 100 / means[time_point][blank_idx][peak] * np.sqrt(
                    (stds[time_point][idx][peak])**2 +
                    ((means[time_point][idx][peak] - means[time_point][blank_idx][peak]) / means[time_point][blank_idx][peak])**2 * (stds[time_point][blank_idx][peak])**2
                )
            means_diff.append(diff)
            stds_diff.append(err)

        axs[channel_idx].bar(x + A*0.1, means_diff, yerr=stds_diff, width=0.1, align='center', alpha=0.5, capsize=5, label=labels[A])
        df2[f'{channel} Intensity'] = means_diff   #####
        df2[f'{channel} Intensity_Std'] = stds_diff   #####

    axs[channel_idx].set_xticks(x)
    axs[channel_idx].set_xticklabels(x)
    axs[channel_idx].legend()

# Save and display the plots
plt.tight_layout()
plt.savefig(f'660-norm-All-Intensity change-t={time_point}.png', dpi=300, bbox_inches='tight')
plt.show()

df2_dropped = df2.drop(columns=['DNA Number'])    #####
combined_df = pd.concat([df, df2_dropped], axis=1)   #####
combined_df.to_excel('Intensity_peakshift_Data_By_Chirall660.xlsx', index=False)  #####



               
                       ###Kinetic Study -Intensity and Peak shift change### 

##slides for Intensity

from pptx import Presentation
from pptx.util import Inches


# itial intensity
df2 = pd.read_csv('DNAClusteringGlucoseSelectivityBefore-660 dec.csv', delimiter=',', header=None)

## Calculate the means and stds

total=60  # ensure total is an integer

# create lists for the means and stds
meansI0 = []
stdsI0 = []

for i in range(20):  # within each time point, step by 1
    indices = [i + j*20 for j in range(3)]
        # ensure all indices are within time point columns
    if all(idx < total for idx in indices):
        meansI0.append(df2.iloc[:, indices].mean(axis=1).tolist())
        stdsI0.append(df2.iloc[:, indices].std(axis=1).tolist())


## combine

meansT=[meansI0]+means
stdsT=[stdsI0]+stds


## plot
time_points = ['I0','1', '3', '5', '7','9','11']
line_styles = ['-', '-', '-', '-', ':', ':', ':', ':']
colors = ['b', 'g', 'r', 'k', 'b', 'g', 'r', 'k']
custom_labels = ["DNA-Glucose 20mM", "DNA-PBS 10 %"]

peaks = [8, 14]
channels = ['(7,5)', '(7,6)']

##slides for Intensity
peaks = [7, 13]
channels = ['(7,5)', '(7,6)']


prs = Presentation()

# Create an empty slide layout (this might differ depending on your PowerPoint version)
slide_layout = prs.slide_layouts[6]

for idx in range(12, 16):
    indices = [idx, idx + 4]
    
    # Prepare plots for Intensity change and Peak position change
    for peaks, ylabel, image_name in zip([[8, 14], [7, 13]], ['Intensity change (%)', 'Peak position (%)'], ['Comparation_660-DNA-', 'Comparation-intensity_660-DNA-']):
        fig, axs = plt.subplots(1, 2, figsize=(8, 3))
        axs = axs.ravel()

        for ax, peak, channel in zip(axs, peaks, channels):
            for i in indices:
                means1 = [meansT[j][i][peak] for j in range(7)]
                stds1 = [stdsT[j][i][peak] for j in range(7)]
                ax.errorbar(time_points, means1, yerr=stds1, label=custom_labels[indices.index(i)], linestyle=line_styles[indices.index(i)], color=colors[indices.index(i)])
            ax.set_xlabel('Time Points (hr)', fontsize=14)
            ax.set_ylabel(ylabel, fontsize=14)
            ax.legend()
            ax.set_title(f'{channel}')

        plt.tight_layout()
        plt.savefig(f'{image_name}{idx}.png', dpi=300, bbox_inches='tight')

    # Add a new slide with the empty layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Define the position and size of the images on the slide
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(3)

    # Insert the images into the slide
    for image_name in ['Comparation_660-DNA-', 'Comparation-intensity_660-DNA-']:
        img_path = f'{image_name}{idx}.png'
        slide.shapes.add_picture(img_path, left, top, width, height)
        top += height  # Adjust top position for next image

# Save the presentation
prs.save('kinetic_study-660-1-none.pptx')

plt.show()



                        ## compare with Intiall intensity ##
                        
                        ## PLot the result  ###              

# Specify the time point and row to calculate the differences
## peak shift
time_point=0
time_pointI0 = 0
time_pointI = 0

peaks = [7,13]
channels = ['(7,5)', '(7,6)']
x_values = np.arange(1, 49)

# Create the subplots
fig, axs = plt.subplots(1, 2, figsize=(13, 3))

for i, (ax, peak, channel) in enumerate(zip(axs, peaks, channels)):
    means_diff = [meansT[time_pointI][col][peak] - meansT[time_pointI0][col][peak] for col in range(48)]
    stds_diff = [np.sqrt(stdsT[time_pointI][col][peak]**2 + stdsT[time_pointI0][col][peak]**2) for col in range(48)]
    
    # Set the current subplot to ax
    plt.sca(ax)
    
    # Specify bar colors based on y-values
    bar_colors = ['g' if y > 0.8 else 'g' if y < -0.8 else 'b' for y in means_diff]

    plt.bar(x_values, means_diff, yerr=stds_diff, align='center', alpha=0.5, capsize=0, 
             color=bar_colors, error_kw={'ecolor': 'black', 'elinewidth': 1})
    plt.title(f'{channel}')
    plt.xlabel('DNA Number',fontsize=12,fontweight='bold')
    plt.ylabel('Peak Shift (nm)',fontsize=12,fontweight='bold')
    plt.ylim(-1, 1)  # set the limits of y-axis
    plt.xticks(x_values)
    plt.xticks(x_values, rotation=90)

# Save and display the subplots
plt.tight_layout()
plt.savefig(f'660-All-Peak shiftCAMPARE 2Io-t={time_point}.png', dpi=300, bbox_inches='tight')
plt.show()



# Intensity change
peaks = [8,14]
channels = ['(7,5)', '(7,6)']
x_values = np.arange(1, 49)

# Create the subplots
fig, axs = plt.subplots(1, 2, figsize=(15, 6))

for i, (ax, peak, channel) in enumerate(zip(axs, peaks, channels)):
    # Calculate the differences in means and standard deviations
    means_diff = [(meansT[time_pointI][col][peak] - meansT[time_pointI0][col][peak]) / meansT[time_pointI0][col][peak] * 100 for col in range(48)]

    # Calculate the propagated errors for standard deviation differences
    stds1 = [stdsT[time_pointI][col][peak] for col in range(48)]
    stds2 = [stdsT[time_pointI0][col][peak] for col in range(48)]
    stds_diff = [np.sqrt(std1**2 + std2**2) / meansT[time_pointI0][col][peak] * 100 for std1, std2, col in zip(stds1, stds2, range(48))]
            # Set the current subplot to ax
    plt.sca(ax)
        # Specify bar colors based on y-values
    bar_colors = ['g' if y > 30 else 'g' if y < -30 else 'b' for y in means_diff]

    plt.bar(x_values, means_diff, yerr=stds_diff, align='center', alpha=0.5, capsize=0, 
             color=bar_colors, error_kw={'ecolor': 'black', 'elinewidth': 1})
    plt.title(f'{channel}')
    plt.xlabel('DNA Number',fontsize=12,fontweight='bold')
    plt.ylabel('Intensity change (%)',fontsize=12,fontweight='bold')
    #plt.ylim(-1, 1)  # set the limits of y-axis
    plt.xticks(x_values)
    plt.xticks(x_values, rotation=90)

# Save and display the subplots
plt.tight_layout()
plt.savefig(f'660-All-Intensity change CAMPARE 2Io-t={time_point}.png', dpi=300, bbox_inches='tight')
plt.show()



        

        
### Normalized Intensity



from pptx import Presentation
from pptx.util import Inches


# itial intensity
df2 = pd.read_csv('DNAClusteringGlucose 660 dec.csv', delimiter=',', header=None)

## Calculate the means and stds

total=288  # ensure total is an integer

# create lists for the means and stds
meansI0 = []
stdsI0 = []

for i in range(96):  # within each time point, step by 1
    indices = [i + j*96 for j in range(3)]
        # ensure all indices are within time point columns
    if all(idx < total for idx in indices):
        meansI0.append(df2.iloc[:, indices].mean(axis=1).tolist())
        stdsI0.append(df2.iloc[:, indices].std(axis=1).tolist())


## combine

meansT=[meansI0]+means
stdsT=[stdsI0]+stds


## plot
time_points = ['I0','2', '4', '6', '8']
line_styles = ['-', '-', '-', '-', ':', ':', ':', ':']
colors = ['b', 'g', 'r', 'k', 'b', 'g', 'r', 'k']
custom_labels = ["DNA-Glucose 20mM", "DNA-PBS 10 %"]

peaks = [8, 14]
channels = ['(7,5)', '(7,6)']

##slides for Intensity
peaks = [7, 13]
channels = ['(7,5)', '(7,6)']


prs = Presentation()
for idx in range(0, 48):
    indices = [idx, idx + 48]
    
    # Prepare plots for Intensity change and Peak position change
    for peaks, ylabel, image_name in zip([[8, 14], [7, 13]], ['Normalized Intensity change', 'Peak position change'], ['Comparation_660-DNA-', 'Comparation-intensity_660-DNA-']):
        fig, axs = plt.subplots(1, 2, figsize=(8, 3))
        axs = axs.ravel()

        for ax, peak, channel in zip(axs, peaks, channels):
            for i in indices:
                # Get the intensity or peak position values
                values = [meansT[j][i][peak] for j in range(5)]
                # Get the standard deviation values
                stds1 = [stdsT[j][i][peak] for j in range(5)]

                # Normalize the intensities by dividing by the initial intensity
                if ylabel == 'Normalized Intensity change':
                    means1 = [value / values[0] for value in values]
                    # Normalize the standard deviations using the error propagation rule
                    stds1_norm = [std / values[0] * means1[i] for i, std in enumerate(stds1)]
                    # Set y-limits for the intensity change plot
                    ax.set_ylim([0, 2])  # adjust the limits as needed
                else:
                    means1 = values
                    stds1_norm = stds1
                
                ax.errorbar(time_points, means1, yerr=stds1_norm, label=custom_labels[indices.index(i)], linestyle=line_styles[indices.index(i)], color=colors[indices.index(i)])
                
            ax.set_xlabel('Time Points (hr)', fontsize=14)
            ax.set_ylabel(ylabel, fontsize=14)
            ax.legend()
            ax.set_title(f'{channel}')

        plt.tight_layout()
        plt.savefig(f'{image_name}{idx}.png', dpi=300, bbox_inches='tight')

    # Add a new slide with the empty layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Define the position and size of the images on the slide
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(3)

    # Insert the images into the slide
    for image_name in ['Comparation_660-DNA-', 'Comparation-intensity_660-DNA-']:
        img_path = f'{image_name}{idx}.png'
        slide.shapes.add_picture(img_path, left, top, width, height)
        top += height  # Adjust top position for next image
# Save the presentation
prs.save('kinetic_study-660-Normalized.pptx')

plt.show()
